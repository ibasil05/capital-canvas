"""
PowerPoint export functionality for financial models.
"""

import io
import os
from pathlib import Path
from typing import Dict, Any, List, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import config

class PowerPointExport:
    """PowerPoint export handler for financial models"""
    
    # Default template path
    TEMPLATE_DIR = Path(__file__).parent.parent.parent / "templates"
    DEFAULT_TEMPLATE = TEMPLATE_DIR / "financial_model_template.pptx"
    
    def __init__(self, model_data: Dict[str, Any], ticker: str, company_name: str):
        """
        Initialize PowerPoint export handler.
        
        Args:
            model_data: Financial model data
            ticker: Company ticker
            company_name: Company name
        """
        self.model_data = model_data
        self.ticker = ticker
        self.company_name = company_name
        
        # Check if template exists, otherwise create a new presentation
        if self.DEFAULT_TEMPLATE.exists():
            self.prs = Presentation(self.DEFAULT_TEMPLATE)
        else:
            self.prs = Presentation()
    
    def generate(self) -> bytes:
        """
        Generate PowerPoint file containing the financial model.
        
        Returns:
            PowerPoint file as bytes
        """
        slide_methods = [
            self._create_title_slide,
            self._create_summary_slide,
            self._create_dcf_valuation_slide,
            self._create_comps_valuation_slide,
            self._create_lbo_analysis_slide,
            self._create_income_statement_slide,
            self._create_balance_sheet_slide,
            self._create_cash_flow_slide,
            self._create_capital_structure_slide,
        ]

        for method in slide_methods:
            try:
                method()
            except Exception as slide_err:
                # Log and continue building remaining slides
                print(f"[PowerPointExport] Skipping slide {method.__name__} due to error: {slide_err}")
        
        # Save to bytes IO
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        
        return output.getvalue()
    
    def _create_title_slide(self):
        """Create the title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = None
        try:
            subtitle = slide.placeholders[1]
        except IndexError:
            # Some templates may not have a second placeholder; fall back to adding a textbox
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
            subtitle = subtitle_box.text_frame
        
        title.text = f"{self.company_name} ({self.ticker})"
        if hasattr(subtitle, "text"):
            subtitle.text = "Financial Model & Valuation Analysis"
        else:
            subtitle.text = "Financial Model & Valuation Analysis"  # TextFrame also supports .text
    
    def _create_summary_slide(self):
        """Create the model summary slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Financial Model Summary"
        
        # Add content as a table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8)
        
        table_rows = 4
        table_cols = 3
        
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(3)
        table.columns[2].width = Inches(3)
        
        # Add table headers
        headers = ["DCF Valuation", "Trading Comps", "LBO Analysis"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Add key metrics
        dcf_data = self.model_data.get("dcf_valuation", {})
        comps_data = self.model_data.get("trading_comps_valuation", {})
        lbo_data = self.model_data.get("lbo_valuation", {})
        
        metrics = [
            ["Enterprise Value", 
             f"${fmt_num(dcf_data.get('enterprise_value'), 1_000_000_000, suffix='B')}", 
             f"${fmt_num(comps_data.get('enterprise_value'), 1_000_000_000, suffix='B')}", 
             f"${fmt_num(lbo_data.get('entry_enterprise_value'), 1_000_000_000, suffix='B')}"],
            ["Equity Value", 
             f"${fmt_num(dcf_data.get('equity_value'), 1_000_000_000, suffix='B')}", 
             f"${fmt_num(comps_data.get('equity_value'), 1_000_000_000, suffix='B')}", 
             f"${fmt_num(lbo_data.get('entry_equity_value'), 1_000_000_000, suffix='B')}"],
            ["Share Price / IRR", 
             f"${fmt_num(dcf_data.get('price_per_share'), precision=2)}", 
             f"${fmt_num(comps_data.get('price_per_share'), precision=2)}", 
             f"{fmt_num(lbo_data.get('equity_irr'), pct=True)}%"]
        ]
        
        for i, row_data in enumerate(metrics):
            metric_name = row_data[0]
            cell = table.cell(i+1, 0)
            cell.text = metric_name
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            
            for j in range(3):
                cell = table.cell(i+1, j)
                if j > 0:  # Skip the metric name cell which we already set
                    cell.text = row_data[j+1]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
        
        # Add a chart showing key metrics
        self._add_summary_chart(slide)
    
    def _add_summary_chart(self, slide):
        """Add a chart to the summary slide"""
        # Get data for the chart
        income_data = self.model_data.get("income_statement", {})
        years = list(range(5))  # Assume 5 years of forecasts
        
        if isinstance(income_data, dict) and "revenue" in income_data:
            revenue = [income_data["revenue"].get(str(year), 0) / 1_000_000_000 for year in years]  # Convert to billions
            ebitda = [income_data["ebitda"].get(str(year), 0) / 1_000_000_000 for year in years]  # Convert to billions
            
            # Create chart
            chart_data = CategoryChartData()
            chart_data.categories = [f'Year {year+1}' for year in years]
            chart_data.add_series('Revenue ($B)', revenue)
            chart_data.add_series('EBITDA ($B)', ebitda)
            
            x, y, cx, cy = Inches(1), Inches(3.5), Inches(8), Inches(3.5)
            try:
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                ).chart
                chart.has_legend = True
                chart.has_title = True
                chart.chart_title.text_frame.text = "Revenue and EBITDA Forecast"
            except Exception as chart_err:
                print(f"[PowerPointExport] Skipping summary chart due to error: {chart_err}")
    
    def _create_dcf_valuation_slide(self):
        """Create the DCF valuation slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Discounted Cash Flow (DCF) Valuation"
        
        # Add DCF key metrics and assumptions
        dcf_data = self.model_data.get("dcf_valuation", {})
        
        # Try to get content placeholder; fallback to new textbox
        try:
            content_tf = slide.placeholders[1].text_frame
        except IndexError:
            content_tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(6)).text_frame
        content = content_tf
        
        p = content.paragraphs[0]
        p.text = "Key Metrics and Assumptions"
        p.font.bold = True
        p.font.size = Pt(18)
        
        metrics = [
            ("Discount Rate (WACC)", f"{fmt_num(dcf_data.get('discount_rate'), pct=True)}%"),
            ("Terminal Growth Rate", f"{fmt_num(dcf_data.get('terminal_growth_rate'), pct=True)}%"),
            ("PV of Forecast Cash Flows", f"${fmt_num(dcf_data.get('pv_forecast_fcf'), 1_000_000_000, suffix='B')}"),
            ("PV of Terminal Value", f"${fmt_num(dcf_data.get('pv_terminal_value'), 1_000_000_000, suffix='B')}"),
            ("Enterprise Value", f"${fmt_num(dcf_data.get('enterprise_value'), 1_000_000_000, suffix='B')}"),
            ("Equity Value", f"${fmt_num(dcf_data.get('equity_value'), 1_000_000_000, suffix='B')}"),
            ("Implied Share Price", f"${fmt_num(dcf_data.get('price_per_share'), precision=2)}")
        ]
        
        for metric, value in metrics:
            p = content.add_paragraph()
            p.text = f"{metric}: {value}"
        
        # Add a chart showing DCF breakdown
        self._add_dcf_chart(slide)
    
    def _add_dcf_chart(self, slide):
        """Add a chart showing DCF value breakdown"""
        dcf_data = self.model_data.get("dcf_valuation", {})
        
        pv_fcf = dcf_data.get("pv_forecast_fcf", 0) / 1_000_000_000  # Convert to billions
        pv_tv = dcf_data.get("pv_terminal_value", 0) / 1_000_000_000  # Convert to billions
        net_debt = dcf_data.get("net_debt", 0) / 1_000_000_000  # Convert to billions
        
        # Create chart
        chart_data = CategoryChartData()
        chart_data.categories = ['PV of FCF', 'PV of Terminal Value', 'Net Debt', 'Equity Value']
        chart_data.add_series('Value ($B)', [pv_fcf, pv_tv, -net_debt, pv_fcf + pv_tv - net_debt])
        
        x, y, cx, cy = Inches(5), Inches(2), Inches(4), Inches(4)
        try:
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = False
            chart.has_title = True
            chart.chart_title.text_frame.text = "DCF Value Bridge"
        except Exception as chart_err:
            print(f"[PowerPointExport] Skipping DCF chart due to error: {chart_err}")
    
    def _create_comps_valuation_slide(self):
        """Create the trading comps valuation slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Trading Comparables Valuation"
        
        # Add trading comps metrics
        comps_data = self.model_data.get("trading_comps_valuation", {})
        
        # Try to get content placeholder; fallback to new textbox
        try:
            content_tf = slide.placeholders[1].text_frame
        except IndexError:
            content_tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(6)).text_frame
        content = content_tf
        
        p = content.paragraphs[0]
        p.text = "Valuation Metrics"
        p.font.bold = True
        p.font.size = Pt(18)
        
        metrics = [
            ("Forward EBITDA", f"${fmt_num(comps_data.get('forward_ebitda'), 1_000_000, suffix='M')}"),
            ("EV/EBITDA Multiple", f"{fmt_num(comps_data.get('ev_to_ebitda'), suffix='x')}"),
            ("EV/Revenue Multiple", f"{fmt_num(comps_data.get('ev_to_revenue'), suffix='x')}"),
            ("P/E Ratio", f"{fmt_num(comps_data.get('price_to_earnings'), suffix='x')}"),
            ("Enterprise Value", f"${fmt_num(comps_data.get('enterprise_value'), 1_000_000_000, suffix='B')}"),
            ("Equity Value", f"${fmt_num(comps_data.get('equity_value'), 1_000_000_000, suffix='B')}"),
            ("Implied Share Price", f"${fmt_num(comps_data.get('price_per_share'), precision=2)}")
        ]
        
        for metric, value in metrics:
            p = content.add_paragraph()
            p.text = f"{metric}: {value}"
        
        # Add trading comps table
        self._add_comps_table(slide)
    
    def _add_comps_table(self, slide):
        """Add a table with trading comparables"""
        # Create a table for trading comps
        left = Inches(5)
        top = Inches(2)
        width = Inches(4)
        height = Inches(4)
        
        peers = self.model_data.get("trading_comps", [])
        if not peers:
            # Generate some example peers if none exist
            peers = [
                {"ticker": "PEER1", "name": "Peer Company 1", "ev_to_ebitda": 7.5, "ev_to_revenue": 1.8, "price_to_earnings": 14.0},
                {"ticker": "PEER2", "name": "Peer Company 2", "ev_to_ebitda": 8.2, "ev_to_revenue": 2.1, "price_to_earnings": 16.5},
                {"ticker": "PEER3", "name": "Peer Company 3", "ev_to_ebitda": 9.0, "ev_to_revenue": 2.3, "price_to_earnings": 18.0}
            ]
        
        table_rows = len(peers) + 2  # Header + median + peers
        table_cols = 4
        
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
        
        # Set headers
        headers = ["Company", "EV/EBITDA", "EV/Revenue", "P/E"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add peer data
        for i, peer in enumerate(peers):
            table.cell(i+1, 0).text = peer.get("ticker", "")
            table.cell(i+1, 1).text = f"{fmt_num(peer.get('ev_to_ebitda'), suffix='x')}"
            table.cell(i+1, 2).text = f"{fmt_num(peer.get('ev_to_revenue'), suffix='x')}"
            table.cell(i+1, 3).text = f"{fmt_num(peer.get('price_to_earnings'), suffix='x')}"
        
        # Add median row
        median_row = len(peers) + 1
        table.cell(median_row, 0).text = "Median"
        
        # Calculate medians
        ev_ebitda_values = [peer.get("ev_to_ebitda", 0) for peer in peers]
        ev_revenue_values = [peer.get("ev_to_revenue", 0) for peer in peers]
        pe_values = [peer.get("price_to_earnings", 0) for peer in peers]
        
        ev_ebitda_median = sorted(ev_ebitda_values)[len(ev_ebitda_values)//2] if ev_ebitda_values else 0
        ev_revenue_median = sorted(ev_revenue_values)[len(ev_revenue_values)//2] if ev_revenue_values else 0
        pe_median = sorted(pe_values)[len(pe_values)//2] if pe_values else 0
        
        table.cell(median_row, 1).text = f"{fmt_num(ev_ebitda_median, suffix='x')}"
        table.cell(median_row, 2).text = f"{fmt_num(ev_revenue_median, suffix='x')}"
        table.cell(median_row, 3).text = f"{fmt_num(pe_median, suffix='x')}"
        
        # Style median row
        for i in range(4):
            cell = table.cell(median_row, i)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
    
    def _create_lbo_analysis_slide(self):
        """Create the LBO analysis slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Leveraged Buyout (LBO) Analysis"
        
        # Add LBO metrics
        lbo_data = self.model_data.get("lbo_valuation", {})
        
        # Try to get content placeholder; fallback to new textbox
        try:
            content_tf = slide.placeholders[1].text_frame
        except IndexError:
            content_tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(6)).text_frame
        content = content_tf
        
        p = content.paragraphs[0]
        p.text = "LBO Analysis Results"
        p.font.bold = True
        p.font.size = Pt(18)
        
        metrics = [
            ("Holding Period", f"{fmt_num(lbo_data.get('holding_period_years'), precision=0)} years"),
            ("Exit Multiple", f"{fmt_num(lbo_data.get('exit_multiple'), suffix='x')}"),
            ("Entry Enterprise Value", f"${fmt_num(lbo_data.get('entry_enterprise_value'), 1_000_000_000, suffix='B')}"),
            ("Entry Equity Value", f"${fmt_num(lbo_data.get('entry_equity_value', lbo_data.get('equity_investment')), 1_000_000_000, suffix='B')}"),
            ("Initial Debt", f"${fmt_num(lbo_data.get('entry_debt', lbo_data.get('debt_investment')), 1_000_000_000, suffix='B')}"),
            ("Exit Enterprise Value", f"${fmt_num(lbo_data.get('exit_enterprise_value'), 1_000_000_000, suffix='B')}"),
            ("Exit Equity Value", f"${fmt_num(lbo_data.get('exit_equity_value'), 1_000_000_000, suffix='B')}"),
            ("Equity IRR", f"{fmt_num(lbo_data.get('equity_irr'), pct=True)}%"),
            ("Cash-on-Cash Multiple", f"{fmt_num(lbo_data.get('cash_on_cash', lbo_data.get('cash_on_cash_multiple')), suffix='x')}"),
            ("Entry Debt/EBITDA", f"{fmt_num(lbo_data.get('entry_debt_to_ebitda'), suffix='x')}"),
            ("Exit Debt/EBITDA", f"{fmt_num(lbo_data.get('exit_debt_to_ebitda'), suffix='x')}"),
        ]
        
        for metric, value in metrics:
            p = content.add_paragraph()
            p.text = f"{metric}: {value}"
        
        # Add a chart showing the LBO returns
        self._add_lbo_chart(slide)
    
    def _add_lbo_chart(self, slide):
        """Add a chart showing LBO returns"""
        lbo_data = self.model_data.get("lbo_valuation", {})
        
        entry_equity = lbo_data.get("entry_equity_value", 0) / 1_000_000_000  # Convert to billions
        exit_equity = lbo_data.get("exit_equity_value", 0) / 1_000_000_000  # Convert to billions
        
        # Create chart
        chart_data = CategoryChartData()
        chart_data.categories = ['Entry', 'Exit']
        chart_data.add_series('Equity Value ($B)', [entry_equity, exit_equity])
        
        x, y, cx, cy = Inches(5), Inches(2), Inches(4), Inches(3)
        try:
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = False
            chart.has_title = True
            chart.chart_title.text_frame.text = "LBO Equity Value Growth"
        except Exception as chart_err:
            print(f"[PowerPointExport] Skipping LBO chart due to error: {chart_err}")
    
    def _create_income_statement_slide(self):
        """Create the income statement slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Income Statement Projections"
        
        # Add income statement table
        income_data = self.model_data.get("income_statement", {})
        # Derive available years from the income_data dictionary keys
        def _collect_years(data_dict):
            years_set = set()
            for inner in data_dict.values():
                if isinstance(inner, dict):
                    years_set.update(inner.keys())
            return sorted(years_set, key=lambda y: int(y))

        years = _collect_years(income_data) or ["0","1","2","3","4","5"]
        
        # Create table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(4)
        
        table_rows = 8  # Selected key metrics
        table_cols = len(years) + 1  # Years + row labels
        
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
        
        # Set headers
        table.cell(0, 0).text = "In millions, USD"
        for i, year in enumerate(years):
            col = i + 1
            if i == 0:
                table.cell(0, col).text = "Historical"
            else:
                table.cell(0, col).text = f"Year {year}"
            
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add key income statement items
        items = [
            ("Revenue", "revenue"),
            ("Gross Profit", "gross_profit"),
            ("EBITDA", "ebitda"),
            ("EBITDA Margin", "ebitda_margin"),
            ("Operating Income", "operating_income"),
            ("Net Income", "net_income"),
            ("EPS", "eps")
        ]
        
        for i, (label, key) in enumerate(items):
            row = i + 1
            table.cell(row, 0).text = label
            
            # Format as percentage for margin items
            is_percentage = "margin" in key
            
            # Add values for each year
            if isinstance(income_data, dict) and key in income_data:
                for j, year in enumerate(years):
                    col = j + 1
                    value = income_data[key].get(str(year), 0)
                    
                    if is_percentage:
                        table.cell(row, col).text = f"{fmt_num(value, pct=True)}%"
                    else:
                        # Format in millions
                        value_in_millions = value / 1_000_000
                        table.cell(row, col).text = f"${fmt_num(value_in_millions)}"
    
    def _create_balance_sheet_slide(self):
        """Create the balance sheet slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Balance Sheet Projections"
        
        # Add balance sheet table
        balance_data = self.model_data.get("balance_sheet", {})
        years = list(range(6))  # Historical + 5 year forecast
        
        # Create table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(4)
        
        items = [
            ("Cash", "cash"),
            ("Accounts Receivable", "accounts_receivable"),
            ("Inventory", "inventory"),
            ("Total Current Assets", "total_current_assets"),
            ("Fixed Assets", "fixed_assets"),
            ("Total Assets", "total_assets"),
            ("Total Current Liabilities", "total_current_liabilities"),
            ("Total Debt", "total_debt"),
            ("Total Equity", "total_equity")
        ]
        
        table_rows = len(items) + 1  # header + metrics
        
        table = slide.shapes.add_table(table_rows, 1, left, top, width, height).table
        
        # Set headers
        table.cell(0, 0).text = "In millions, USD"
        for i, (label, key) in enumerate(items):
            table.cell(i+1, 0).text = label
            
            # Add values for each year
            if isinstance(balance_data, dict) and key in balance_data:
                for j, year in enumerate(years):
                    col = j + 1
                    value = balance_data[key].get(str(year), 0)
                    
                    # Format in millions
                    value_in_millions = value / 1_000_000
                    table.cell(i+1, col).text = f"${fmt_num(value_in_millions)}"
    
    def _create_cash_flow_slide(self):
        """Create the cash flow statement slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Cash Flow Projections"
        
        # Add cash flow table
        cash_flow_data = self.model_data.get("cash_flow", {})
        years = list(range(6))  # Historical + 5 year forecast
        
        # Create table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(4)
        
        items = [
            ("Net Income", "net_income"),
            ("Depreciation & Amortization", "depreciation"),
            ("Change in Working Capital", "change_in_working_capital"),
            ("Cash Flow from Operations", "operating_cash_flow"),
            ("Capital Expenditures", "capex"),
            ("Cash Flow from Investing", "investing_cash_flow"),
            ("Free Cash Flow", "free_cash_flow")
        ]
        
        table_rows = len(items) + 1  # header + metrics
        
        table = slide.shapes.add_table(table_rows, 1, left, top, width, height).table
        
        # Set headers
        table.cell(0, 0).text = "In millions, USD"
        for i, (label, key) in enumerate(items):
            table.cell(i+1, 0).text = label
            
            # Add values for each year
            if isinstance(cash_flow_data, dict) and key in cash_flow_data:
                for j, year in enumerate(years):
                    col = j + 1
                    value = cash_flow_data[key].get(str(year), 0)
                    
                    # Format in millions
                    value_in_millions = value / 1_000_000
                    table.cell(i+1, col).text = f"${fmt_num(value_in_millions)}"
    
    def _create_capital_structure_slide(self):
        """Create the capital structure analysis slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Capital Structure Analysis"
        
        # Add capital structure grid
        cap_structure_data = self.model_data.get("capital_structure_grid", [])
        
        # Create table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(3)
        
        table_rows = min(6, len(cap_structure_data) + 1)  # Header + up to 5 scenarios
        table_cols = 6  # Selected key metrics
        
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
        
        # Set headers
        headers = ["Debt/EBITDA", "Debt/Capital", "WACC", "Credit Rating", "Equity IRR", "Share Price"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add capital structure scenarios
        for i, scenario in enumerate(cap_structure_data[:5]):  # Limit to 5 scenarios
            row = i + 1
            
            table.cell(row, 0).text = f"{fmt_num(scenario.get('debt_to_ebitda'), suffix='x')}"
            table.cell(row, 1).text = f"{fmt_num(scenario.get('debt_to_capital'), pct=True)}%"
            table.cell(row, 2).text = f"{fmt_num(scenario.get('wacc'), pct=True)}%"
            table.cell(row, 3).text = str(scenario.get('credit_rating', ""))
            table.cell(row, 4).text = f"{fmt_num(scenario.get('equity_irr'), pct=True)}%"
            table.cell(row, 5).text = f"${fmt_num(scenario.get('share_price'), precision=2)}"
        
        # Add a chart
        self._add_capital_structure_chart(slide)
    
    def _add_capital_structure_chart(self, slide):
        """Add a chart showing the capital structure tradeoffs"""
        cap_structure_data = self.model_data.get("capital_structure_grid", [])
        
        if not cap_structure_data:
            return
        
        # Extract data for chart
        debt_to_ebitda = [scenario.get("debt_to_ebitda", 0) for scenario in cap_structure_data[:5]]
        wacc = [scenario.get("wacc", 0) * 100 for scenario in cap_structure_data[:5]]  # Convert to percentage
        
        # Create chart
        chart_data = CategoryChartData()
        chart_data.categories = [f"{fmt_num(d, suffix='x')}" for d in debt_to_ebitda]
        chart_data.add_series('WACC (%)', wacc)
        
        x, y, cx, cy = Inches(2), Inches(5), Inches(6), Inches(3)
        try:
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = True
            chart.has_title = True
            chart.chart_title.text_frame.text = "WACC vs. Leverage"
            chart.value_axis.has_major_gridlines = True
        except Exception as chart_err:
            print(f"[PowerPointExport] Skipping capital-structure chart due to error: {chart_err}")

def _safe_float(val, default=0.0):
    try:
        return float(val)
    except (TypeError, ValueError):
        return default

def fmt_num(value, scale=1, precision=1, suffix='', pct=False):
    """Safely format a number with scaling and suffix. Returns 'N/A' if not convertible."""
    try:
        num = float(value) / scale
        if pct:
            num *= 100
        return f"{num:.{precision}f}{suffix}"
    except (TypeError, ValueError):
        return "N/A" 