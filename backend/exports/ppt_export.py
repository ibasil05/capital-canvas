from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from typing import Dict, List, Any
from datetime import datetime
from .powerpoint import fmt_num  # Import safe number formatter
from .powerpoint import PowerPointExport  # Full-featured exporter

# Slide Layouts (assuming standard layouts)
TITLE_SLIDE_LAYOUT = 0
TITLE_AND_CONTENT_LAYOUT = 1
SECTION_HEADER_LAYOUT = 2 
BLANK_LAYOUT = 5
CONTENT_WITH_CAPTION_LAYOUT = 7 # Example, might vary

# Common Styling
FONT_NAME = "Inter"
TITLE_FONT_SIZE = Pt(32)
SUBTITLE_FONT_SIZE = Pt(18)
BODY_FONT_SIZE = Pt(12)
SMALL_FONT_SIZE = Pt(10)
BLACK_COLOR = RGBColor(0, 0, 0)
GREY_COLOR = RGBColor(107, 111, 118) # #6B6F76 from PRD

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if title: title.text = title_text
    
    subtitle = slide.placeholders[1] # Placeholder index for subtitle can vary
    if subtitle: subtitle.text = subtitle_text
    return slide

def add_content_slide(prs, title_text, layout_idx=TITLE_AND_CONTENT_LAYOUT):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    if title_shape: title_shape.text = title_text
    # Return slide and the main content placeholder (index typically 1 for content layouts)
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1 or shape.name.startswith("Content Placeholder") or shape.name.startswith("Text Placeholder"): # Common indices/names
            content_placeholder = shape
            break
    return slide, content_placeholder

async def generate_ppt_export(model_results_data: Dict[str, Any]) -> bytes:
    """High-level wrapper that produces a banker-grade 10-slide deck.
    Delegates to backend.exports.powerpoint.PowerPointExport which contains the
    full logic for summary, valuation, statements, heat-map, etc.
    """

    # Extract high-level metadata from the results (fallbacks in case keys missing)
    ticker = model_results_data.get("ticker") or model_results_data.get("symbol") or "COMPANY"
    company_name = model_results_data.get("company_name") or model_results_data.get("name") or ticker

    try:
        exporter = PowerPointExport(model_data=model_results_data,
                                    ticker=ticker.upper(),
                                    company_name=company_name)
        return exporter.generate()
    except Exception as e:
        # As a fallback (and to avoid entirely breaking export), fall back to the
        # lightweight placeholder deck implemented earlier.
        print(f"[ppt_export] Fallback to minimal deck after error in PowerPointExport: {e}")
        prs = Presentation()
        add_title_slide(prs, f"{company_name} ({ticker})", "Financial Model & Valuation Overview")
        slide, _ = add_content_slide(prs, "Export Error")
        slide.shapes.placeholders[1].text_frame.text = str(e)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read() 